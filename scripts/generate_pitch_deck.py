"""Generate FLOWRA suite pitch deck — 4 files:
  - pitch_deck_pointers.pdf     (10 slides, lean)
  - pitch_deck_pointers.pptx    (10 slides, lean)
  - pitch_deck_detailed.pdf     (16 slides + write-up appendix)
  - pitch_deck_detailed.pptx    (16 slides, full walkthrough)

Brand:
  navy   #0F1B4C   primary
  blue   #2563EB   accent / headline
  amber  #f59e0b   CTA / highlight
  paper  #FFFFFF   bg
  grey   #64748B   body
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from reportlab.lib import colors as rc
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import mm
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle, Image,
)

OUT = Path("/app/frontend/public/pitch")
OUT.mkdir(parents=True, exist_ok=True)

POSTERS = Path("/app/frontend/public/posters")
SALESMAN_POSTER = POSTERS / "flowra-salesman-poster.jpg"
DISPATCH_POSTER = POSTERS / "flowra-dispatch-poster.jpg"

# ───── BRAND ──────────────────────────────────────────────────────────────
NAVY = "#0F1B4C"
BLUE = "#2563EB"
AMBER = "#f59e0b"
PAPER = "#FFFFFF"
GREY = "#64748B"
SOFT = "#F0F4FF"

# ───── CONTENT ────────────────────────────────────────────────────────────
COMPANY = "Jodidar India"
CITY = "Raipur, Chhattisgarh, India"
PRODUCT = "FLOWRA Insights"
CONCLAVE = "Manthan X Founders"
FOUNDER = "Ankit Sarawgi"
TEAM = "Punit · Kritika"
TAGLINE = "Organize. Automate. Accelerate."

# ───── SLIDE DECK STRUCTURE (10 pointer | 16 detailed) ───────────────────
# Each slide: dict(title, bullets, sub, note)
#   bullets = list of strings (one-liners for pointer deck)
#   sub = sub-heading line under the title
#   note = (detailed deck) longer paragraph appended below bullets
# For both decks we share the same 16 source slides; pointer deck just
# picks 10 of them and trims bullets/notes.

ALL_SLIDES = [
    {
        "tag": "cover",
        "title": "FLOWRA Insights",
        "sub": "Organize. Automate. Accelerate.",
        "bullets": [
            "An AI-first operating system for Indian SMEs running Tally / Busy",
            f"Made by {COMPANY}, {CITY}",
            f"Founder: {FOUNDER}  ·  Team: {TEAM}",
            f"Pitched at: {CONCLAVE}",
        ],
        "note": "",
    },
    {
        "tag": "problem",
        "title": "The Problem",
        "sub": "10 million Indian SMEs run their books on Tally / Busy — and are stuck there.",
        "bullets": [
            "Tally / Busy were built in the 1990s — desktop-only, single-machine, no analytics layer",
            "Owners can't see live sales, stock, dispatch or salesman performance unless they sit at the office machine",
            "Field salesmen capture orders on paper / WhatsApp → 2–4 day data lag",
            "Dispatch & receivables are manual — papers move from desk to desk",
            "Hindi-English bilingual + low-bandwidth realities are ignored by most SaaS",
        ],
        "note": (
            "Indian MSMEs contribute ~30% of GDP. Almost every one of them uses Tally ERP 9 / Tally Prime "
            "or Busy as their book of record — but those systems were never designed for the smartphone "
            "era. The result is a structural blind-spot: the owner knows what was in his books last "
            "Saturday, not what his salesman or godown did this morning."
        ),
    },
    {
        "tag": "opportunity",
        "title": "The Opportunity",
        "sub": "₹15,000 Cr SME-software TAM in India by 2027.",
        "bullets": [
            "TAM — 10 M+ active SMEs on Tally / Busy across India",
            "SAM — 1.5 M SMEs with ≥3 employees and a smartphone-using owner",
            "SOM — 50 K early-adopter SMEs in Tier-2 & Tier-3 cities (Year 1 wedge)",
            "Tailwind — Digital India, ONDC, GST 3.0, Account Aggregator → SME digitization is now policy",
        ],
        "note": (
            "We are NOT competing with Tally — we sit ON TOP of it. The Tally Sync Agent reads vouchers, "
            "stock, masters, ledgers via Tally's own ODBC / TDL XML protocol. Zero workflow change for "
            "the accountant. Owner gets the cloud + AI + mobile layer."
        ),
    },
    {
        "tag": "solution",
        "title": "The Solution — FLOWRA Insights",
        "sub": "One cloud layer that sees everything Tally / Busy sees — plus AI on top.",
        "bullets": [
            "Lightweight Windows agent reads Tally / Busy in real time (no manual export)",
            "Multi-company, multi-branch, multi-FY ready out of the box",
            "Mobile-first dashboards for the owner, salesman, dispatch & CA",
            "AI Insights powered by Gemini & GPT-5 — explains numbers in plain Hindi/English",
            "100% Indian data residency — Mumbai region MongoDB Atlas",
        ],
        "note": (
            "Plug-and-play: install agent → log in → 30-second sync of last 3 years of data. The agent "
            "auto-updates from our manifest. Every API endpoint enforces tenant isolation; no two "
            "customer's data can ever cross."
        ),
    },
    {
        "tag": "modules",
        "title": "The Suite — 8 Integrated Modules",
        "sub": "Each module ships value on day 1; together they replace 6–8 disjoint tools.",
        "bullets": [
            "1. Sales & Receivables — live aging, top customers, collection nudges",
            "2. Inventory & Stock — ABC, reorder, multi-group filters, PO generation",
            "3. Salesman App — beat plans, daily check-ins, order & payment capture",
            "4. Dispatch Terminal — Kanban board, LR / Bilty upload, porter assignment",
            "5. CA Corner — TDS, GST, balance-sheet reconciliation, journal nudges",
            "6. AI Insights — natural-language Q&A on your own books",
            "7. FLOWRA Tasks — task & workflow assignment across roles",
            "8. FLOWRA Loyalty — built-in customer-rewards engine for SMEs",
        ],
        "note": (
            "Modules share a unified data model — a sales voucher captured by the salesman flows "
            "instantly into Sales reports, Dispatch board, CA reconciliation and the Loyalty points "
            "ledger. One source of truth, eight surfaces."
        ),
    },
    {
        "tag": "differentiators",
        "title": "Differentiators",
        "sub": "Why not just build it in Tally / Excel / Zoho?",
        "bullets": [
            "Tally-native — no double-entry, no CSV imports, no workflow change",
            "Bilingual UI — English + Hindi, Tier-2/Tier-3 ready",
            "Multi-tenant SaaS — full data isolation, GST-compliant audit trail",
            "AI layer — GPT-5 + Gemini, explains receivables / inventory in plain language",
            "FLOWRA Tasks & Loyalty — value-adds Tally / Zoho simply don't have",
            "Made-in-India price points — fraction of Salesforce / Zoho One",
        ],
        "note": (
            "Competitors fall into 3 camps: (a) Tally add-ons that still need a desktop; (b) global "
            "SaaS like Zoho / Salesforce that ignore Tally; (c) standalone ERPs that demand migration. "
            "FLOWRA sits in the white-space — Tally-native, cloud-mobile, AI-aware, and priced for India."
        ),
    },
    {
        "tag": "tasks",
        "title": "FLOWRA Tasks — Roles, Reminders, Accountability",
        "sub": "A built-in operations layer for the owner and the team.",
        "bullets": [
            "Assign tasks to salesman, dispatch, CA, or any employee — by name or by role",
            "Recurring tasks (weekly stock count, monthly GSTR, daily bank reco)",
            "Tally-data-aware triggers — auto-create a task when receivables cross a threshold",
            "Mobile push + WhatsApp reminders (roadmap)",
            "Full audit trail — who did what, when, with proof attachment",
        ],
        "note": (
            "Tasks isn't a generic to-do list. Every task can be linked to a Tally voucher, customer, "
            "or item — so closing a task automatically updates the relevant FLOWRA report. Think "
            "Asana, but it actually knows your books."
        ),
    },
    {
        "tag": "loyalty",
        "title": "FLOWRA Loyalty — Reward Programs, Built-In",
        "sub": "Turn one-time buyers into repeat customers without a separate tool.",
        "bullets": [
            "Points engine — auto-credit on every Tally sales voucher",
            "Tiers — Silver / Gold / Platinum with tenant-defined rules",
            "Redemption — discount coupon or cash-back against next invoice",
            "Customer-facing WhatsApp / SMS / email card (roadmap)",
            "Owner dashboard — top loyal customers, churn-risk customers, lifetime value",
        ],
        "note": (
            "Indian SMEs lose ~25-30% of their customer base every year because they have no "
            "structured retention play. FLOWRA Loyalty fixes that — without forcing the owner to "
            "subscribe to a separate platform like LoyaltyXpert, Capillary or Easyrewardz."
        ),
    },
    {
        "tag": "how",
        "title": "How It Works",
        "sub": "5 steps from sign-up to first AI insight — under 10 minutes.",
        "bullets": [
            "1. Owner signs up at flowralive.in (free 7-day trial)",
            "2. Downloads the FLOWRA Tally Sync Agent (10 MB Windows .exe)",
            "3. Agent reads Tally vouchers + masters in real time via official TDL",
            "4. Data lands in encrypted Mumbai-region MongoDB Atlas (multi-tenant)",
            "5. Dashboards + AI Insights + mobile apps light up instantly",
        ],
        "note": (
            "The Sync Agent is the moat. It auto-discovers companies, handles Tally Prime 7.0's new "
            "AlterID protocol, supports multi-FY, and auto-updates itself from a signed manifest. "
            "Customers don't touch their accountant's setup."
        ),
    },
    {
        "tag": "traction",
        "title": "Where We Are",
        "sub": "Live with pilot customers — feedback-driven product.",
        "bullets": [
            "Pilot customers across automobile, electrical, textile & lubricant verticals",
            "100K+ Tally vouchers processed in pilots",
            "v9.8.28 desktop agent — Tally Prime 7.0 compatible, hot-update channel live",
            "111 product iterations shipped; 101/101 automated tests green",
            "Bilingual marketing posters + landing page ready (flowralive.in)",
        ],
        "note": (
            "We deliberately picked pilot customers across 4 verticals to harden multi-tenant, "
            "multi-company and bilingual edge cases before scale-up. Their daily feedback drives "
            "the public roadmap on the website."
        ),
    },
    {
        "tag": "business_model",
        "title": "Business Model",
        "sub": "Subscription SaaS · per company · per user · annual upfront.",
        "bullets": [
            "Tiered subscription: Starter / Growth / Enterprise",
            "Per-additional-user pricing for salesman & dispatch logins",
            "Add-on revenue: WhatsApp credits, AI Calling minutes, GST filing assistance",
            "Pricing — currently pilot test going on; commercial rollout post-pilot",
            "Sticky integration → low expected churn once a tenant onboards",
        ],
        "note": (
            "Once a tenant maps its Tally companies, customers, salesmen and loyalty rules into "
            "FLOWRA, the cost of switching to another tool is high. We expect <5% annual churn "
            "based on comparable Indian SaaS benchmarks (Vyapar, RazorpayX Payroll)."
        ),
    },
    {
        "tag": "gtm",
        "title": "Go-to-Market",
        "sub": "Three concentric circles — CA → Tally Partner → SME owner.",
        "bullets": [
            "1. Partner-led — Chartered Accountants + Tally Partners channel (revenue share)",
            "2. Self-serve — Hindi/English landing page, 7-day free trial, no credit card",
            "3. Community-led — Tier-2/3 owner communities, WhatsApp groups, local conclaves",
            "4. Content — bilingual demo videos, GST/TDS explainers, AI-insight walkthroughs",
            "5. Referral engine — built-in referral codes for existing users (live in product)",
        ],
        "note": (
            "Indian SMEs trust their CA more than any marketer. We turn the CA into FLOWRA's first "
            "champion — they get a Corner of their own (CA Corner module), and a revenue share when "
            "their clients subscribe."
        ),
    },
    {
        "tag": "tech",
        "title": "Tech Stack",
        "sub": "Modern, multi-tenant, India-first.",
        "bullets": [
            "Frontend — React 19 + Tailwind, mobile-first, bilingual",
            "Backend — Python FastAPI, async, hot-reload, modular routes",
            "DB — MongoDB Atlas (Mumbai region) with full tenant isolation",
            "Desktop Agent — Windows Python + PyInstaller, signed .exe, auto-update",
            "AI — Gemini 3 + GPT-5 + nano-banana via Emergent integrations",
            "Hosting — DigitalOcean droplet, Nginx + Supervisor, encrypted at rest",
        ],
        "note": (
            "We chose Python + React + MongoDB so a small founding team can ship daily. Multi-tenant "
            "isolation is enforced at every query; no application code can ever leak data across "
            "tenants. All secrets are in environment variables, not in code."
        ),
    },
    {
        "tag": "team",
        "title": "The Team",
        "sub": "Small, founder-led, ship-every-day culture.",
        "bullets": [
            f"{FOUNDER} — Founder, Product & Strategy",
            "Punit — Engineering & Tally protocol",
            "Kritika — UX, Brand & Customer Success",
            f"Headquartered in {CITY}",
            "AI-augmented engineering — daily build, daily test, daily ship",
        ],
        "note": (
            "Our biggest moat is operational velocity — we ship and verify against live customer "
            "Tally data every week. A larger competitor would take a quarter to do what we ship in a sprint."
        ),
    },
    {
        "tag": "roadmap",
        "title": "12-Month Roadmap",
        "sub": "From pilot to 500+ paying SMEs.",
        "bullets": [
            "Q1 — WhatsApp automation (AiSensy / Meta Cloud) + AI Calling bot for recovery",
            "Q2 — GST Portal integration + FLOWRA Loyalty v1 GA + Google Drive document store",
            "Q3 — Hindi UI complete + ONDC connector + Open Account-Aggregator pull",
            "Q4 — Iframe / API for ERP partners + India-first AI co-pilot",
            "Always-on — daily product polish driven by pilot customer feedback",
        ],
        "note": (
            "Every roadmap item is already specced and prioritised in our PRD.md (versioned in repo). "
            "We do not run on guesswork — every quarter's deliverables are visible to pilot customers."
        ),
    },
    {
        "tag": "closing",
        "title": "Thank You",
        "sub": "Let's build the operating system for Indian SMEs together.",
        "bullets": [
            "Product — flowralive.in",
            "Founder — Ankit Sarawgi  ·  ankit@jodidarindia.in",
            f"Company — {COMPANY}, {CITY}",
            f"Pitched at — {CONCLAVE}",
            "Questions welcome — and demo available on request.",
        ],
        "note": "",
    },
]

POINTER_TAGS = [
    "cover", "problem", "opportunity", "solution", "modules",
    "differentiators", "how", "traction", "business_model", "closing",
]


# ────────────────────────────────────────────────────────────────────────────
# POWERPOINT
# ────────────────────────────────────────────────────────────────────────────
def _rgb(hex_: str) -> RGBColor:
    return RGBColor.from_string(hex_.lstrip("#"))


def add_bg(slide, prs, color_hex):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(color_hex)
    bg.shadow.inherit = False
    return bg


def add_accent_bar(slide, prs):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(BLUE)


def add_footer(slide, prs, text):
    tb = slide.shapes.add_textbox(Inches(0.6), prs.slide_height - Inches(0.45),
                                  prs.slide_width - Inches(1.2), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9)
    r.font.color.rgb = _rgb(GREY)
    r.font.name = "Calibri"


def add_text(slide, x, y, w, h, text, *, size=24, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = _rgb(color)
    r.font.name = font
    return tb


def render_pptx(slides, out_path, deck_label):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, s in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        add_bg(slide, prs, PAPER)
        add_accent_bar(slide, prs)
        is_cover = s["tag"] == "cover"
        is_closing = s["tag"] == "closing"

        if is_cover:
            # Hero layout
            add_text(slide, Inches(0.8), Inches(1.4), Inches(8), Inches(0.5),
                     CONCLAVE.upper(), size=14, bold=True, color=BLUE)
            add_text(slide, Inches(0.8), Inches(1.85), Inches(11), Inches(1.6),
                     s["title"], size=64, bold=True, color=NAVY)
            add_text(slide, Inches(0.8), Inches(3.5), Inches(11), Inches(0.7),
                     s["sub"], size=24, color=BLUE)
            # Bullets without bullets — meta strip
            y = Inches(4.5)
            for b in s["bullets"]:
                add_text(slide, Inches(0.8), y, Inches(11), Inches(0.4),
                         "·  " + b, size=14, color=GREY)
                y += Inches(0.4)
            # Embed poster on right
            try:
                if SALESMAN_POSTER.exists():
                    slide.shapes.add_picture(str(SALESMAN_POSTER), Inches(9.5),
                                             Inches(1.4), Inches(3.2), Inches(3.2))
            except Exception:
                pass
        elif is_closing:
            add_text(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(1.4),
                     s["title"], size=64, bold=True, color=NAVY)
            add_text(slide, Inches(0.8), Inches(3.5), Inches(11), Inches(0.7),
                     s["sub"], size=22, color=BLUE)
            y = Inches(4.3)
            for b in s["bullets"]:
                add_text(slide, Inches(0.8), y, Inches(11), Inches(0.45),
                         "·  " + b, size=16, color=NAVY)
                y += Inches(0.45)
            try:
                if DISPATCH_POSTER.exists():
                    slide.shapes.add_picture(str(DISPATCH_POSTER), Inches(9.5),
                                             Inches(1.8), Inches(3.2), Inches(3.2))
            except Exception:
                pass
        else:
            # Standard content layout
            add_text(slide, Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.5),
                     f"{deck_label}  ·  Slide {idx + 1} / {len(slides)}",
                     size=10, bold=True, color=BLUE)
            add_text(slide, Inches(0.8), Inches(1.05), Inches(11.5), Inches(0.95),
                     s["title"], size=34, bold=True, color=NAVY)
            add_text(slide, Inches(0.8), Inches(2.05), Inches(11.5), Inches(0.55),
                     s["sub"], size=16, color=BLUE)
            # Bullets
            y = Inches(2.95)
            for b in s["bullets"]:
                # Accent dot
                dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), y + Inches(0.12),
                                             Inches(0.12), Inches(0.12))
                dot.line.fill.background()
                dot.fill.solid()
                dot.fill.fore_color.rgb = _rgb(BLUE)
                add_text(slide, Inches(1.15), y, Inches(11.0), Inches(0.45),
                         b, size=15, color=NAVY)
                y += Inches(0.5)
            # Note block (only for detailed deck — pointer deck strips notes)
            if s.get("note") and deck_label.lower().startswith("detailed"):
                box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8),
                                             prs.slide_height - Inches(1.85),
                                             Inches(11.7), Inches(1.2))
                box.line.fill.background()
                box.fill.solid()
                box.fill.fore_color.rgb = _rgb(SOFT)
                add_text(slide, Inches(1.0), prs.slide_height - Inches(1.75),
                         Inches(11.3), Inches(1.0), s["note"], size=11, color=NAVY)

        add_footer(slide, prs, f"{COMPANY} · {PRODUCT} · {TAGLINE}")

    prs.save(out_path)


# ────────────────────────────────────────────────────────────────────────────
# PDF
# ────────────────────────────────────────────────────────────────────────────
def _styles():
    s = getSampleStyleSheet()
    return {
        "h1": ParagraphStyle("h1", parent=s["Title"], fontSize=42, leading=46,
                             textColor=rc.HexColor(NAVY), spaceAfter=8,
                             alignment=0),
        "sub": ParagraphStyle("sub", parent=s["Normal"], fontSize=16, leading=20,
                              textColor=rc.HexColor(BLUE), spaceAfter=12),
        "bullet": ParagraphStyle("bullet", parent=s["Normal"], fontSize=13, leading=18,
                                 textColor=rc.HexColor(NAVY), spaceAfter=4,
                                 leftIndent=14, bulletIndent=2),
        "note": ParagraphStyle("note", parent=s["Normal"], fontSize=10, leading=14,
                               textColor=rc.HexColor(NAVY), spaceBefore=8, spaceAfter=4,
                               backColor=rc.HexColor(SOFT), borderPadding=8),
        "footer": ParagraphStyle("footer", parent=s["Normal"], fontSize=8, leading=10,
                                 textColor=rc.HexColor(GREY), alignment=1),
        "writeup_h": ParagraphStyle("writeup_h", parent=s["Heading2"], fontSize=18,
                                    textColor=rc.HexColor(NAVY), spaceBefore=12, spaceAfter=6),
        "writeup_p": ParagraphStyle("writeup_p", parent=s["Normal"], fontSize=11, leading=16,
                                    textColor=rc.HexColor(NAVY), spaceAfter=8, alignment=4),
        "meta": ParagraphStyle("meta", parent=s["Normal"], fontSize=10, leading=14,
                               textColor=rc.HexColor(BLUE), spaceAfter=4),
        "small": ParagraphStyle("small", parent=s["Normal"], fontSize=9, leading=12,
                                textColor=rc.HexColor(GREY)),
    }


def _draw_slide_chrome(canvas, doc, deck_label, total_pages_func):
    canvas.saveState()
    # Top accent bar
    canvas.setFillColor(rc.HexColor(BLUE))
    canvas.rect(0, doc.pagesize[1] - 6 * mm, doc.pagesize[0], 6 * mm, fill=1, stroke=0)
    # Footer
    canvas.setFillColor(rc.HexColor(GREY))
    canvas.setFont("Helvetica", 8)
    canvas.drawString(15 * mm, 8 * mm, f"{COMPANY} · {PRODUCT} · {TAGLINE}")
    canvas.drawRightString(doc.pagesize[0] - 15 * mm, 8 * mm,
                           f"{deck_label} · Slide {canvas.getPageNumber()}")
    canvas.restoreState()


def _slide_flow(s, st, deck_label, is_detailed):
    flow = []
    if s["tag"] == "cover":
        flow.append(Spacer(1, 16 * mm))
        flow.append(Paragraph(CONCLAVE.upper(), st["meta"]))
        flow.append(Spacer(1, 4 * mm))
        flow.append(Paragraph(s["title"], st["h1"]))
        flow.append(Paragraph(s["sub"], st["sub"]))
        flow.append(Spacer(1, 10 * mm))
        for b in s["bullets"]:
            flow.append(Paragraph(f"·  {b}", st["bullet"]))
        flow.append(Spacer(1, 8 * mm))
        if SALESMAN_POSTER.exists():
            try:
                flow.append(Image(str(SALESMAN_POSTER), width=70 * mm, height=70 * mm))
            except Exception:
                pass
    elif s["tag"] == "closing":
        flow.append(Spacer(1, 24 * mm))
        flow.append(Paragraph(s["title"], st["h1"]))
        flow.append(Paragraph(s["sub"], st["sub"]))
        flow.append(Spacer(1, 6 * mm))
        for b in s["bullets"]:
            flow.append(Paragraph(f"·  {b}", st["bullet"]))
    else:
        flow.append(Paragraph(s["title"], st["h1"]))
        flow.append(Paragraph(s["sub"], st["sub"]))
        flow.append(Spacer(1, 4 * mm))
        for b in s["bullets"]:
            flow.append(Paragraph(f"<b>·</b>  {b}", st["bullet"]))
        if is_detailed and s.get("note"):
            flow.append(Paragraph(s["note"], st["note"]))
    flow.append(PageBreak())
    return flow


def render_pdf(slides, out_path, deck_label, *, is_detailed, with_writeup=False):
    st = _styles()
    doc = SimpleDocTemplate(
        out_path,
        pagesize=landscape(A4),
        leftMargin=18 * mm, rightMargin=18 * mm,
        topMargin=20 * mm, bottomMargin=15 * mm,
        title=f"{PRODUCT} — {deck_label}",
        author=COMPANY,
    )

    flow = []
    for s in slides:
        flow.extend(_slide_flow(s, st, deck_label, is_detailed))

    if with_writeup:
        flow.extend(_writeup_appendix(st))

    def _chrome(canvas, _doc):
        _draw_slide_chrome(canvas, _doc, deck_label, lambda: len(slides))

    doc.build(flow, onFirstPage=_chrome, onLaterPages=_chrome)


# ────────────────────────────────────────────────────────────────────────────
# WRITE-UP APPENDIX (Detailed PDF only)
# ────────────────────────────────────────────────────────────────────────────
WRITEUP_SECTIONS = [
    ("Project Title",
     f"<b>{PRODUCT}</b> — An AI-first cloud operating system for Indian "
     f"micro, small and medium enterprises running Tally / Busy ERP."),
    ("Promoter & Company",
     f"<b>{COMPANY}</b>, head-quartered in {CITY}.  Founder: <b>{FOUNDER}</b>. "
     f"Core team: {TEAM}. The company is a founder-led, bootstrapped product "
     "venture building software for the unserved long-tail of Indian SMEs."),
    ("Problem Statement",
     "Over ten million Indian SMEs run their day-to-day accounting on Tally "
     "ERP 9, Tally Prime, or Busy. These tools, while excellent for compliance, "
     "are essentially single-machine desktop systems built in the 1990s. They "
     "do not expose live data to the smartphone, do not offer analytics, and "
     "cannot natively serve the field salesman, the warehouse, the chartered "
     "accountant or the customer. The result is delayed decisions, paper-based "
     "operations, and a structural drag on SME productivity."),
    ("Proposed Solution",
     "FLOWRA Insights is a lightweight cloud SaaS that sits on top of Tally / "
     "Busy without disturbing the accountant's existing workflow. A small "
     "Windows agent reads Tally vouchers, stock and masters in real time via "
     "the official TDL/ODBC protocol and pushes them to a multi-tenant cloud "
     "where the owner, salesman, dispatch, and CA each have a purpose-built "
     "mobile-first dashboard. An AI layer powered by Gemini 3 and GPT-5 turns "
     "raw numbers into plain-language insights and recommendations."),
    ("Product Modules",
     "The suite ships eight integrated modules — Sales &amp; Receivables, "
     "Inventory &amp; Stock, Salesman App with Beat Run, Dispatch Terminal "
     "Kanban, CA Corner, AI Insights, <b>FLOWRA Tasks</b> (task &amp; "
     "workflow assignment with Tally-data-aware triggers), and <b>FLOWRA "
     "Loyalty</b> (a built-in points / tier customer-rewards engine for "
     "SMEs). Every module shares a unified tenant-isolated data model, so "
     "a sales voucher flows instantly into reports, dispatch, CA reconciliation "
     "and loyalty points."),
    ("Differentiation",
     "Most Indian SaaS competitors either (a) sit inside Tally as desktop "
     "add-ons, (b) ignore Tally and demand migration, or (c) target only "
     "specific functions like invoicing. FLOWRA is the only product that is "
     "Tally-native, cloud-mobile, AI-aware, bilingual, and priced for Tier-2 "
     "/ Tier-3 SMEs all at once."),
    ("Market Opportunity",
     "The Indian SME-software TAM is projected at ~₹15,000 crore by 2027. "
     "Within this, the addressable market for cloud + AI overlays on existing "
     "Tally / Busy installations is ~₹3,500 crore, driven by Digital India, "
     "ONDC adoption, GST 3.0 and the Account Aggregator framework. Our serviceable "
     "obtainable market in Year 1 is 50,000 early-adopter SMEs in Tier-2 / "
     "Tier-3 cities, where founder Ankit Sarawgi has strong network advantages."),
    ("Business Model",
     "Subscription SaaS — billed per company per user, annual upfront. "
     "Tiered packaging (Starter / Growth / Enterprise) with add-on revenue "
     "from WhatsApp messaging credits, AI calling minutes and GST filing "
     "assistance. Pricing is currently in <b>pilot test</b> with feedback from "
     "live customers; commercial rollout is planned at the end of the pilot phase."),
    ("Go-to-Market",
     "Three concentric circles: (i) Chartered Accountants and Tally Partners "
     "as channel partners with revenue share — the CA Corner module gives "
     "them a daily reason to recommend FLOWRA to their clients; (ii) self-serve "
     "via a bilingual landing page with 7-day free trial and no credit card; "
     "(iii) community-led growth through Tier-2 / Tier-3 SME WhatsApp groups "
     "and local startup conclaves such as Manthan X Founders."),
    ("Technology",
     "React 19 + Tailwind frontend, FastAPI (Python 3.11) backend, MongoDB "
     "Atlas in the Mumbai region, Python Windows desktop agent (signed .exe, "
     "auto-updating). AI is delivered via Emergent integrations with Gemini 3, "
     "GPT-5 and nano-banana image generation. All secrets live in environment "
     "variables; every endpoint enforces tenant isolation at the query layer."),
    ("Team",
     f"{FOUNDER} (Founder — product, strategy, customer development). "
     "Punit (Engineering — Tally protocol, agent, backend). "
     "Kritika (UX, brand, customer success). The team operates on a "
     "daily-build / daily-ship cadence with AI-augmented engineering."),
    ("Traction",
     "111 product iterations shipped to date. 101 of 101 automated tests "
     "green across the test suite. Live pilots across automobile, electrical, "
     "textile and lubricant verticals. 100,000+ Tally vouchers processed in "
     "pilot data. v9.8.28 of the desktop agent in production, with hot-update "
     "channel and Tally Prime 7.0 compatibility."),
    ("Roadmap (Next 12 Months)",
     "Q1 — WhatsApp Automation and AI Calling Bot for receivables recovery. "
     "Q2 — GST Portal integration, FLOWRA Loyalty v1 GA, Google Drive document "
     "storage. Q3 — Hindi-complete UI, ONDC connector, Account Aggregator pull. "
     "Q4 — Partner API / iframe, India-first AI co-pilot. Always-on: daily "
     "product polish driven by pilot customer feedback."),
    ("Risks & Mitigations",
     "(a) Tally protocol drift — mitigated by the modular, hot-updatable agent "
     "and a comprehensive integration test suite. (b) Multi-tenant data leak — "
     "mitigated by query-layer enforcement and a 100% green tenant-isolation "
     "test pack. (c) Adoption inertia — mitigated by zero-workflow-change "
     "installation, 7-day free trial and CA-partner-led pull. "),
    ("Ask",
     "We are pitching at the <b>Manthan X Founders</b> conclave to invite "
     "mentorship, partnerships, and the wider grant / investor ecosystem to "
     "join us in building the operating system for Indian SMEs."),
]


def _writeup_appendix(st):
    flow = [PageBreak(),
            Paragraph(f"{PRODUCT} — Project Write-up", st["h1"]),
            Paragraph(f"For submission at {CONCLAVE}", st["sub"]),
            Spacer(1, 6 * mm)]
    for heading, body in WRITEUP_SECTIONS:
        flow.append(Paragraph(heading, st["writeup_h"]))
        flow.append(Paragraph(body, st["writeup_p"]))
    flow.append(Spacer(1, 6 * mm))
    flow.append(Paragraph(
        f"<i>Prepared by {COMPANY}, {CITY}. © FLOWRA Insights. {TAGLINE}</i>",
        st["small"],
    ))
    return flow


# ────────────────────────────────────────────────────────────────────────────
def main():
    pointers = [s for s in ALL_SLIDES if s["tag"] in POINTER_TAGS]
    detailed = ALL_SLIDES

    render_pptx(pointers, str(OUT / "flowra_pitch_deck_pointers.pptx"), "Pointers")
    print(f"✓ {OUT / 'flowra_pitch_deck_pointers.pptx'}")
    render_pptx(detailed, str(OUT / "flowra_pitch_deck_detailed.pptx"), "Detailed")
    print(f"✓ {OUT / 'flowra_pitch_deck_detailed.pptx'}")

    render_pdf(pointers, str(OUT / "flowra_pitch_deck_pointers.pdf"),
               "Pointers", is_detailed=False, with_writeup=False)
    print(f"✓ {OUT / 'flowra_pitch_deck_pointers.pdf'}")
    render_pdf(detailed, str(OUT / "flowra_pitch_deck_detailed.pdf"),
               "Detailed", is_detailed=True, with_writeup=True)
    print(f"✓ {OUT / 'flowra_pitch_deck_detailed.pdf'}")


if __name__ == "__main__":
    main()
