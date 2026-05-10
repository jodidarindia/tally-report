"""Render FLOWRA_BUSINESS_PROPOSAL.md → PDF + PPTX.

Outputs:
  /app/docs/FLOWRA_BUSINESS_PROPOSAL.pdf   (polished investor-grade PDF)
  /app/docs/FLOWRA_BUSINESS_PROPOSAL.pptx  (19-slide pitch deck w/ charts)

The PDF flows the entire proposal with custom styling.
The PPTX is a hand-crafted 19-slide deck — not an auto-conversion — built
for live presentation in front of an investor or a strategic partner.
Includes 3 matplotlib-rendered visual chart slides.
"""
import os
import re
from pathlib import Path
import markdown
from weasyprint import HTML, CSS
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import matplotlib
matplotlib.use("Agg")  # headless backend — no display required
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
import numpy as np

DOCS = Path(__file__).resolve().parent
SRC = DOCS / "FLOWRA_BUSINESS_PROPOSAL.md"
PDF = DOCS / "FLOWRA_BUSINESS_PROPOSAL.pdf"
PPTX = DOCS / "FLOWRA_BUSINESS_PROPOSAL.pptx"
CHARTS_DIR = DOCS / "_charts"
CHARTS_DIR.mkdir(exist_ok=True)


# ─────────────────────── PDF ─────────────────────────────────────────────
PDF_CSS = """
@page {
  size: A4;
  margin: 22mm 18mm 22mm 18mm;
  @top-left   { content: "FLOWRA Insights · Business Proposal";
                font-size: 9pt; color: #64748B; font-family: 'Inter', sans-serif; }
  @top-right  { content: "v1.0 · Feb 2026";
                font-size: 9pt; color: #64748B; font-family: 'Inter', sans-serif; }
  @bottom-right { content: counter(page) " / " counter(pages);
                  font-size: 9pt; color: #94A3B8; font-family: 'Inter', sans-serif; }
  @bottom-left  { content: "CONFIDENTIAL";
                  font-size: 9pt; color: #DC2626; font-family: 'Inter', sans-serif;
                  letter-spacing: 1pt; }
}
@page :first {
  margin: 0;
  @top-left  { content: none; } @top-right { content: none; }
  @bottom-right { content: none; } @bottom-left  { content: none; }
}

html, body {
  font-family: 'Inter', 'Segoe UI', 'Helvetica Neue', Arial, sans-serif;
  font-size: 10.5pt; line-height: 1.55; color: #1E293B;
}

.cover {
  page-break-after: always; height: 100vh;
  background: linear-gradient(135deg, #0F172A 0%, #2563EB 60%, #10B981 130%);
  color: #FFFFFF; padding: 70mm 22mm 22mm 22mm; position: relative;
}
.cover .brand {
  font-size: 11pt; letter-spacing: 6pt; color: #93C5FD;
  margin-bottom: 22mm; text-transform: uppercase;
}
.cover h1 {
  font-size: 44pt; font-weight: 800; line-height: 1.05;
  margin: 0 0 8mm 0; color: #FFFFFF; border-bottom: none; padding: 0;
}
.cover .subtitle {
  font-size: 16pt; color: #DBEAFE; font-weight: 400;
  margin-bottom: 24mm; max-width: 140mm;
}
.cover .meta {
  position: absolute; bottom: 22mm; left: 22mm; right: 22mm;
  border-top: 1px solid #334155; padding-top: 10mm;
  display: flex; justify-content: space-between;
  font-size: 9.5pt; color: #CBD5E1;
}
.cover .meta strong { color: #FFFFFF; font-weight: 600; }
.cover .stamp {
  position: absolute; top: 22mm; right: 22mm;
  border: 2pt solid #DC2626; color: #DC2626;
  padding: 2mm 6mm; font-size: 10pt; font-weight: 700;
  letter-spacing: 1pt; transform: rotate(-6deg);
}

h1 {
  font-size: 22pt; font-weight: 700; color: #0F172A;
  margin: 14mm 0 4mm 0; padding-bottom: 3mm;
  border-bottom: 2px solid #2563EB;
  page-break-before: always; page-break-after: avoid;
}
h1:first-of-type { page-break-before: auto; }
h2 { font-size: 14pt; font-weight: 700; color: #1E293B;
     margin: 9mm 0 3mm 0; page-break-after: avoid; }
h3 { font-size: 11.5pt; font-weight: 600; color: #2563EB;
     margin: 6mm 0 2mm 0; page-break-after: avoid; }
h4 { font-size: 10.5pt; font-weight: 600; color: #475569; margin: 5mm 0 2mm 0; }

p { margin: 0 0 3mm 0; orphans: 3; widows: 3; }
strong { color: #0F172A; font-weight: 600; }
ul, ol { margin: 2mm 0 4mm 0; padding-left: 6mm; }
li { margin-bottom: 1.5mm; }
code {
  font-family: 'JetBrains Mono', 'Consolas', monospace; font-size: 9pt;
  background: #F1F5F9; color: #0F172A; padding: 1pt 4pt; border-radius: 3pt;
}
pre {
  background: #0F172A; color: #E2E8F0;
  font-family: 'JetBrains Mono', 'Consolas', monospace;
  font-size: 8.5pt; line-height: 1.5;
  padding: 4mm 5mm; border-radius: 4pt; margin: 3mm 0 5mm 0;
  overflow-x: auto; page-break-inside: avoid;
}
pre code { background: transparent; color: inherit; padding: 0; }

table {
  width: 100%; border-collapse: collapse; margin: 3mm 0 5mm 0;
  font-size: 9.5pt; page-break-inside: avoid;
}
th {
  background: #1E3A8A; color: #FFFFFF; font-weight: 600; text-align: left;
  padding: 2.5mm 3mm; font-size: 9pt;
  text-transform: uppercase; letter-spacing: 0.3pt;
}
td { padding: 2.2mm 3mm; border-bottom: 1px solid #E2E8F0; vertical-align: top; }
tr:nth-child(even) td { background: #F8FAFC; }

blockquote {
  border-left: 3pt solid #2563EB; background: #F0F9FF;
  padding: 3mm 5mm; margin: 4mm 0; color: #1E40AF; font-style: italic;
}
blockquote p { margin: 0; }
hr { border: none; border-top: 1px solid #CBD5E1; margin: 6mm 0; }
table, pre, blockquote { page-break-inside: avoid; }
"""

COVER_HTML = """
<div class="cover">
  <div class="stamp">CONFIDENTIAL</div>
  <div class="brand">FLOWRA · INSIGHTS</div>
  <h1>Business Proposal</h1>
  <div class="subtitle">
    Tally-native business intelligence for India's ₹3,450 Cr SME analytics
    market. Two-year plan to ₹7.5 Cr ARR and Series A.
  </div>
  <div class="meta">
    <div><strong>Version 1.0</strong><br/>February 2026</div>
    <div style="text-align:right">For investor &amp; partner review<br/>
      <strong>Not for redistribution</strong></div>
  </div>
</div>
"""


def render_pdf():
    md = SRC.read_text(encoding="utf-8")
    html_body = markdown.markdown(
        md, extensions=["fenced_code", "tables", "toc", "sane_lists", "attr_list"],
    )
    html_body = re.sub(r"\[\s\]", "☐", html_body)
    html_body = re.sub(r"\[x\]", "☑", html_body)
    full = ('<!doctype html><html><head><meta charset="utf-8"></head><body>'
            + COVER_HTML + html_body + "</body></html>")
    HTML(string=full, base_url=str(DOCS)).write_pdf(
        str(PDF), stylesheets=[CSS(string=PDF_CSS)],
    )
    print(f"Wrote {PDF}  ({PDF.stat().st_size // 1024} KB)")


# ─────────────────────── PPTX ────────────────────────────────────────────
# Brand palette
NAVY    = RGBColor(0x0F, 0x17, 0x2A)
BLUE    = RGBColor(0x25, 0x63, 0xEB)
LIGHT   = RGBColor(0xF8, 0xFA, 0xFC)
GREEN   = RGBColor(0x10, 0xB9, 0x81)
SLATE   = RGBColor(0x47, 0x55, 0x69)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
RED     = RGBColor(0xDC, 0x26, 0x26)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)


def _add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def _add_text(slide, x, y, w, h, text, *, size=14, bold=False,
              color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = "Calibri"
    r.font.color.rgb = color
    return tb


def _add_bullets(slide, x, y, w, h, bullets, *, size=14, color=SLATE):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = "•  " + b
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"


def _slide_chrome(prs, slide, title, subtitle=None, *, page_n=None, total=19):
    """Common header + footer for non-cover slides.
    If page_n is None, infer from len(prs.slides) — the slide must already
    be added before this is called (which is the convention here).
    """
    if page_n is None:
        page_n = len(prs.slides)
    # Top accent bar
    _add_rect(slide, 0, 0, prs.slide_width, Inches(0.05), BLUE)
    # Title block
    _add_text(slide, Inches(0.5), Inches(0.32), Inches(12.5), Inches(0.7),
              title, size=26, bold=True, color=NAVY)
    if subtitle:
        _add_text(slide, Inches(0.5), Inches(1.02), Inches(12.5), Inches(0.4),
                  subtitle, size=13, color=SLATE)
    # Footer
    _add_text(slide, Inches(0.5), Inches(7.0), Inches(6), Inches(0.3),
              "FLOWRA Insights · Business Proposal · v1.0 · Feb 2026",
              size=9, color=MUTED)
    if page_n is not None:
        _add_text(slide, Inches(11.5), Inches(7.0), Inches(1.8), Inches(0.3),
                  f"{page_n} / {total}", size=9, color=MUTED,
                  align=PP_ALIGN.RIGHT)


def _table(slide, x, y, w, h, rows, *, header=True, font_size=11):
    n_rows = len(rows); n_cols = len(rows[0])
    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.margin_left = cell.margin_right = Pt(8)
            cell.margin_top = cell.margin_bottom = Pt(4)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(font_size)
            r.font.name = "Calibri"
            if header and i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
                r.font.color.rgb = WHITE
                r.font.bold = True
            else:
                if i % 2 == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT
                else:
                    cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
                r.font.color.rgb = NAVY
    return tbl


# ─── matplotlib brand palette + chart helpers ──────────────────────────
PLT_NAVY   = "#0F172A"
PLT_BLUE   = "#2563EB"
PLT_LIGHT  = "#DBEAFE"
PLT_GREEN  = "#10B981"
PLT_AMBER  = "#F59E0B"
PLT_RED    = "#DC2626"
PLT_SLATE  = "#475569"
PLT_MUTED  = "#94A3B8"


def _set_brand_style(ax, *, grid_y=True):
    ax.set_facecolor("#FFFFFF")
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    for spine in ("left", "bottom"):
        ax.spines[spine].set_color(PLT_MUTED)
        ax.spines[spine].set_linewidth(0.8)
    ax.tick_params(colors=PLT_SLATE, labelsize=10)
    if grid_y:
        ax.yaxis.grid(True, color="#E2E8F0", linewidth=0.7, zorder=0)
        ax.set_axisbelow(True)


def chart_funnel():
    """Conversion funnel: 10k impressions → 22 paying tenants per ₹50k spend."""
    stages = [
        ("Meta-ad impressions",    10000, PLT_NAVY),
        ("Lead-form submissions",    400, PLT_BLUE),
        ("Qualified",                200, "#3B82F6"),
        ("Demo calls booked",        120, "#60A5FA"),
        ("Free trials",               72, PLT_AMBER),
        ("Paying tenants",            22, PLT_GREEN),
    ]
    values = [s[1] for s in stages]

    fig, ax = plt.subplots(figsize=(11, 5.5), dpi=160)
    fig.patch.set_facecolor("#FFFFFF")

    max_v = max(values)
    y_positions = np.arange(len(stages))[::-1]
    bar_height = 0.7
    for i, (label, val, color) in enumerate(stages):
        width = val / max_v
        y = y_positions[i]
        # Centred bar (looks like a funnel)
        ax.add_patch(FancyBboxPatch(
            (0.5 - width / 2, y - bar_height / 2),
            width, bar_height,
            boxstyle="round,pad=0,rounding_size=0.04",
            linewidth=0, facecolor=color, alpha=0.95,
        ))
        # Stage label inside the bar (right side)
        ax.text(0.5, y, f"{label}    {val:,}",
                ha="center", va="center",
                fontsize=11, color="white", fontweight="bold")
        # Drop-off % beside the bar
        if i > 0:
            prev = values[i - 1]
            rate = val / prev * 100
            ax.text(0.5 + width / 2 + 0.02, y,
                    f"{rate:.0f}%", va="center",
                    fontsize=10, color=PLT_SLATE)

    ax.set_xlim(-0.05, 1.05)
    ax.set_ylim(-0.5, len(stages) - 0.5)
    ax.axis("off")
    ax.set_title(
        "Acquisition funnel — every ₹50k Meta spend yields ~22 paying tenants",
        fontsize=14, color=PLT_NAVY, fontweight="bold", loc="left", pad=12,
    )
    fig.text(0.02, 0.02,
             "Per cohort. 22 tenants × ₹2,499 = ₹54,978 new MRR.  Payback @78% margin: 3.5 months.",
             fontsize=10, color=PLT_SLATE)
    out = CHARTS_DIR / "funnel.png"
    fig.savefig(out, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def chart_growth():
    """Tenant + MRR growth curve M1 → M24 with sensitivity bands."""
    months = np.array([1, 3, 6, 9, 12, 15, 18, 21, 24])
    base   = np.array([15, 95, 280, 580, 1000, 1370, 1750, 2120, 2500])
    bear   = (base * 0.6).astype(int)
    bull   = (base * 1.4).astype(int)

    fig, ax = plt.subplots(figsize=(11, 5.5), dpi=160)
    fig.patch.set_facecolor("#FFFFFF")
    _set_brand_style(ax)

    # Sensitivity envelope
    ax.fill_between(months, bear, bull, color=PLT_BLUE, alpha=0.10,
                    label="Bear ↔ bull sensitivity (0.6× — 1.4×)")
    # Bear & bull lines (thin)
    ax.plot(months, bear, color=PLT_AMBER, linewidth=1.2, linestyle="--",
            label="Bear: 1,500 tenants @ M24")
    ax.plot(months, bull, color=PLT_GREEN, linewidth=1.2, linestyle="--",
            label="Bull: 3,500 tenants @ M24")
    # Base line (thick)
    ax.plot(months, base, color=PLT_BLUE, linewidth=3.0, marker="o",
            markersize=7, label="Base plan: 2,500 tenants @ M24",
            markerfacecolor="white", markeredgewidth=2)

    # Annotate base milestones
    annotations = [(12, 1000, "M12: 1,000\n₹3 Cr ARR"),
                   (24, 2500, "M24: 2,500\n₹7.5 Cr ARR")]
    for mx, my, txt in annotations:
        ax.annotate(txt, xy=(mx, my),
                    xytext=(mx - 2.5, my + 350),
                    fontsize=10, color=PLT_NAVY, fontweight="bold",
                    arrowprops=dict(arrowstyle="->", color=PLT_NAVY, lw=1))

    ax.set_xlabel("Month", color=PLT_SLATE, fontsize=11)
    ax.set_ylabel("Cumulative paying tenants", color=PLT_SLATE, fontsize=11)
    ax.set_xticks([1, 3, 6, 9, 12, 15, 18, 21, 24])
    ax.set_xlim(0, 25)
    ax.set_ylim(0, 4000)

    # Secondary y-axis: MRR
    ax2 = ax.twinx()
    ax2.spines["top"].set_visible(False)
    ax2.spines["right"].set_color(PLT_MUTED)
    ax2.set_ylabel("MRR (₹ Lakh / month)", color=PLT_BLUE, fontsize=11)
    ax2.set_ylim(0, 4000 * 0.025)        # 1 tenant ≈ ₹0.025 L MRR (₹2,499)
    ax2.tick_params(colors=PLT_BLUE, labelsize=10)

    ax.set_title(
        "Tenant growth · 24-month plan with bear/bull envelope",
        fontsize=14, color=PLT_NAVY, fontweight="bold", loc="left", pad=14,
    )
    ax.legend(loc="upper left", frameon=False, fontsize=9.5,
              labelcolor=PLT_NAVY)

    out = CHARTS_DIR / "growth.png"
    fig.savefig(out, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def chart_unit_econ():
    """CAC payback + LTV/CAC visual — compares us vs SaaS benchmark."""
    fig, axes = plt.subplots(1, 2, figsize=(11, 5.0), dpi=160,
                             gridspec_kw={"width_ratios": [1.05, 1]})
    fig.patch.set_facecolor("#FFFFFF")

    # Left: CAC payback (months) — us vs benchmark
    ax = axes[0]
    _set_brand_style(ax)
    cats = ["FLOWRA\nInsights", "SaaS\nbenchmark"]
    vals = [3.5, 12.0]
    cols = [PLT_BLUE, PLT_MUTED]
    bars = ax.bar(cats, vals, color=cols, width=0.5, zorder=3)
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 0.4,
                f"{v} mo", ha="center", color=PLT_NAVY,
                fontsize=14, fontweight="bold")
    ax.set_ylim(0, 14)
    ax.set_ylabel("CAC payback (months)", color=PLT_SLATE, fontsize=11)
    ax.set_title("Payback in 3.5 months", fontsize=12.5,
                 color=PLT_NAVY, fontweight="bold", loc="left", pad=8)

    # Right: LTV/CAC ratio comparison
    ax = axes[1]
    _set_brand_style(ax, grid_y=False)
    cats2 = ["FLOWRA", "Healthy SaaS\n(industry)", "Struggling\n(industry)"]
    ratios = [13, 4, 1.5]
    cols2 = [PLT_GREEN, PLT_BLUE, PLT_RED]
    bars = ax.barh(cats2[::-1], ratios[::-1], color=cols2[::-1],
                   height=0.55, zorder=3)
    for b, v in zip(bars, ratios[::-1]):
        ax.text(b.get_width() + 0.3, b.get_y() + b.get_height() / 2,
                f"{v}×", va="center",
                color=PLT_NAVY, fontsize=14, fontweight="bold")
    ax.set_xlim(0, 16)
    ax.set_xlabel("LTV / CAC ratio", color=PLT_SLATE, fontsize=11)
    ax.xaxis.grid(True, color="#E2E8F0", linewidth=0.7)
    ax.set_axisbelow(True)
    ax.set_title("13× LTV/CAC — exceptional even by SaaS standards",
                 fontsize=12.5, color=PLT_NAVY, fontweight="bold", loc="left", pad=8)

    fig.suptitle(
        "Unit economics — why this scales profitably",
        fontsize=15, color=PLT_NAVY, fontweight="bold", x=0.05, ha="left", y=0.99,
    )
    fig.tight_layout(rect=[0, 0, 1, 0.95])
    out = CHARTS_DIR / "unit_econ.png"
    fig.savefig(out, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def render_all_charts():
    return {
        "funnel":     chart_funnel(),
        "growth":     chart_growth(),
        "unit_econ":  chart_unit_econ(),
    }


# ─── Slide builders ─────────────────────────────────────────────────────
def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Full-bleed gradient (simulated with two stacked rectangles)
    _add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    _add_rect(s, 0, Inches(4.2), prs.slide_width, Inches(3.3), BLUE)
    # Brand
    _add_text(s, Inches(0.6), Inches(0.6), Inches(8), Inches(0.4),
              "FLOWRA · INSIGHTS", size=12, bold=True, color=RGBColor(0x93, 0xC5, 0xFD))
    # Title
    _add_text(s, Inches(0.6), Inches(2.3), Inches(12), Inches(1.6),
              "The Tally-native\nbusiness intelligence platform.",
              size=44, bold=True, color=WHITE)
    # Subtitle
    _add_text(s, Inches(0.6), Inches(5.1), Inches(12), Inches(0.6),
              "Two-year plan to ₹7.5 Cr ARR · ₹4 Cr seed · Feb 2026",
              size=18, color=RGBColor(0xDB, 0xEA, 0xFE))
    # Stamp
    stamp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(11.0), Inches(0.5), Inches(2.0), Inches(0.5))
    stamp.fill.solid(); stamp.fill.fore_color.rgb = WHITE
    stamp.line.color.rgb = RED
    tf = stamp.text_frame; tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "CONFIDENTIAL"
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = RED
    # Footer line
    _add_text(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.3),
              "For investor & partner review · Not for redistribution",
              size=10, color=MUTED)


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(prs, s, "The Problem",
                  "5-second questions taking 30 minutes inside Tally.")
    rows = [
        ["Owner question", "Today (Tally only)"],
        ["How much is overdue from my top 10 customers?",
         "4-step report → Excel → manual sort"],
        ["Where is my salesman? What did he sell yesterday?",
         "Phone call. WhatsApp screenshot."],
        ["Which 20 SKUs are dead inventory?", "Custom report from CA, billed ₹2,000"],
        ["Did the dispatch boy actually deliver invoice 5421?",
         "Phone call to the godown"],
        ["How does this month compare to same month last year?",
         "Two reports, side by side, on paper"],
    ]
    _table(s, Inches(0.5), Inches(1.7), Inches(12.5), Inches(4.0), rows)
    _add_text(s, Inches(0.5), Inches(6.2), Inches(12.5), Inches(0.7),
              "FLOWRA reduces each of these to a tap.",
              size=18, bold=True, color=BLUE)


def slide_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(prs, s, "The Product",
                  "18 modules already in production. Read-only Tally agent.")
    cols = [
        ("Visibility", ["Owner Dashboard", "Sales Analytics", "Inventory Pareto",
                        "AI Insights (GPT-5.2)"]),
        ("Operations", ["CRM with FIFO ageing", "Salesman Beat Plans",
                        "Dispatch Kanban Terminal", "Sync History"]),
        ("Compliance", ["CA Corner: P&L, BS, Cash Flow", "GST-ready exports",
                        "Audit logs", "Daily backups"]),
    ]
    x = Inches(0.5)
    for title, items in cols:
        _add_rect(s, x, Inches(1.7), Inches(4.2), Inches(4.5), LIGHT, line=MUTED)
        _add_rect(s, x, Inches(1.7), Inches(4.2), Inches(0.5), BLUE)
        _add_text(s, x, Inches(1.78), Inches(4.2), Inches(0.4), title,
                  size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_bullets(s, x + Inches(0.2), Inches(2.4), Inches(3.9), Inches(3.6),
                     items, size=13, color=NAVY)
        x += Inches(4.3)
    _add_text(s, Inches(0.5), Inches(6.4), Inches(12), Inches(0.5),
              "The Tally Sync Agent is our moat — 6+ months of XML edge cases competitors must repeat.",
              size=12, color=SLATE, bold=True)


def slide_market(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(prs, s, "Market Size",
                  "11.5 lakh addressable Indian SMEs · ₹3,450 Cr TAM")
    rows = [
        ["Segment", "Tally + Busy seats", "Fit %", "Addressable"],
        ["Manufacturing SMEs (₹5–500 Cr)", "4.5 L", "70%", "3.2 L"],
        ["Distribution / FMCG", "6.5 L", "80%", "5.2 L"],
        ["Retail chains", "5.0 L", "40%", "2.0 L"],
        ["Pharma / Chemicals", "1.5 L", "75%", "1.1 L"],
        ["TOTAL", "17.5 L", "—", "11.5 L"],
    ]
    _table(s, Inches(0.5), Inches(1.7), Inches(12.5), Inches(3.5), rows)
    # Headline boxes
    boxes = [
        ("0.05% capture", "575 tenants", "₹1.7 Cr ARR", LIGHT, NAVY),
        ("0.5% capture", "5,750 tenants", "₹17 Cr ARR", BLUE, WHITE),
        ("2% capture", "23,000 tenants", "₹68 Cr ARR", GREEN, WHITE),
    ]
    x = Inches(0.5)
    for label, sub1, sub2, fill, fg in boxes:
        _add_rect(s, x, Inches(5.5), Inches(4.1), Inches(1.4), fill, line=MUTED)
        _add_text(s, x, Inches(5.6), Inches(4.1), Inches(0.45), label,
                  size=13, bold=True, color=fg, align=PP_ALIGN.CENTER)
        _add_text(s, x, Inches(6.05), Inches(4.1), Inches(0.4), sub1,
                  size=22, bold=True, color=fg, align=PP_ALIGN.CENTER)
        _add_text(s, x, Inches(6.5), Inches(4.1), Inches(0.4), sub2,
                  size=12, color=fg, align=PP_ALIGN.CENTER)
        x += Inches(4.3)


def slide_competition(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(prs, s, "Competitive Landscape",
                  "Nobody else reads from Tally without asking customer to migrate.")
    rows = [
        ["Player", "Approach", "Where they win", "Where we win"],
        ["Zoho Books", "Cloud accounting", "All-in-one", "Customer must migrate off Tally"],
        ["Vyapar / Marg", "Tally alternative", "Mobile billing", "Replaces Tally — non-starter for 95%"],
        ["Khatabook", "Khata only", "Free tier", "No analytics, inventory, dispatch"],
        ["Tally on Cloud", "Lift-and-shift", "Same UI as Tally", "Still no dashboards or mobile"],
        ["Power BI / Tableau", "BI tools", "Powerful", "Needs IT team; ₹50k+ to set up"],
    ]
    _table(s, Inches(0.5), Inches(1.7), Inches(12.5), Inches(4.0), rows, font_size=11)
    _add_text(s, Inches(0.5), Inches(6.4), Inches(12.5), Inches(0.5),
              "Our wedge: zero-migration. Owner gets dashboards in 20 min. CA stays on Tally.",
              size=13, bold=True, color=BLUE)


def slide_business_model(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(prs, s, "Pricing & Unit Economics",
                  "Blended ARPU ₹2,499/mo · LTV/CAC 13×")
    # Pricing table
    rows = [
        ["Plan", "₹/month", "Companies", "Employees", "Modules"],
        ["Starter", "999", "1", "2", "Dashboard + Sales + Inventory + Sync"],
        ["Professional", "2,499", "3", "5", "+ CRM + Analytics"],
        ["Enterprise", "3,799", "10", "20", "+ Salesman + Dispatch + AI + CA"],
    ]
    _table(s, Inches(0.5), Inches(1.7), Inches(12.5), Inches(2.0), rows)
    # KPIs
    kpis = [
        ("ARPU", "₹2,499", "/mo blended"),
        ("Gross margin", "78%", "post all infra costs"),
        ("CAC", "₹3,200", "blended Y1"),
        ("Payback", "3.5 mo", "industry: 12 mo"),
        ("LTV", "₹85,000", "5-yr, 2% churn"),
        ("LTV/CAC", "13×", "industry: 3-5×"),
    ]
    x = Inches(0.5); y = Inches(4.0)
    for i, (label, val, sub) in enumerate(kpis):
        col = i % 3; row = i // 3
        cx = Inches(0.5) + Inches(4.3) * col
        cy = Inches(4.0) + Inches(1.4) * row
        _add_rect(s, cx, cy, Inches(4.1), Inches(1.2), LIGHT, line=MUTED)
        _add_text(s, cx, cy + Inches(0.05), Inches(4.1), Inches(0.3),
                  label, size=11, color=SLATE, align=PP_ALIGN.CENTER)
        _add_text(s, cx, cy + Inches(0.35), Inches(4.1), Inches(0.5),
                  val, size=26, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        _add_text(s, cx, cy + Inches(0.85), Inches(4.1), Inches(0.3),
                  sub, size=10, color=MUTED, align=PP_ALIGN.CENTER)


def slide_gtm(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(prs, s, "Go-to-Market",
                  "Beachhead: distribution sector in MH/GJ/KA.")
    rows = [
        ["Channel", "% new tenants", "CAC (₹)", "Notes"],
        ["Meta Ads (FB + IG lead forms)", "40%", "4,000", "Geo: MH/GJ/KA"],
        ["WhatsApp organic + referral", "25%", "500", "Founder net + early-customer ref"],
        ["YouTube content (Tally tutorials)", "15%", "1,200", "SEO long-tail"],
        ["CA partner programme", "10%", "3,500", "Year 2 trigger"],
        ["Trade shows + Tally events", "5%", "8,000", "High-trust"],
        ["Outbound (cold WA)", "5%", "6,000", "Closing inbound demos"],
        ["BLENDED", "100%", "3,200", "Drops to ₹2,400 by M18"],
    ]
    _table(s, Inches(0.5), Inches(1.7), Inches(12.5), Inches(4.4), rows, font_size=11)
    _add_text(s, Inches(0.5), Inches(6.3), Inches(12.5), Inches(0.5),
              "Funnel: 10,000 impressions → 22 paying tenants → ₹54,978 new MRR / ₹50k spend.",
              size=12, color=SLATE, bold=True)


def slide_traction(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(prs, s, "Traction & Roadmap",
                  "Already shipped. What's next.")
    # Two columns
    _add_rect(s, Inches(0.5), Inches(1.7), Inches(6.1), Inches(5), LIGHT, line=MUTED)
    _add_text(s, Inches(0.7), Inches(1.85), Inches(5.8), Inches(0.4),
              "Already shipped (Feb 2026)", size=15, bold=True, color=GREEN)
    _add_bullets(s, Inches(0.7), Inches(2.3), Inches(5.8), Inches(4.5),
                 ["18 production modules + 120+ regression tests",
                  "Tally Sync Agent v9.8.9 with Day-Book LVD fallback",
                  "Multi-tenant + 6 roles (admin, dispatch, salesman, employee, super_admin, flowra_staff)",
                  "Windows GUI + system-tray + auto-start build kit",
                  "AI Insights (GPT-5.2 expense narratives)",
                  "Production Operations Playbook (18 pages)"],
                 size=12, color=NAVY)
    _add_rect(s, Inches(6.85), Inches(1.7), Inches(6.1), Inches(5), LIGHT, line=MUTED)
    _add_text(s, Inches(7.05), Inches(1.85), Inches(5.8), Inches(0.4),
              "Next 12 months", size=15, bold=True, color=BLUE)
    _add_bullets(s, Inches(7.05), Inches(2.3), Inches(5.8), Inches(4.5),
                 ["Q1: Atlas + DO + CI/CD + Sentry (production hardening)",
                  "Q2: GST Portal + WhatsApp BSP + Audit log CSV",
                  "Q3: Mobile app (React Native), advanced beat optimiser",
                  "Q4: CA partner programme + white-label (alpha)",
                  "Code-signed installer + auto-update channel"],
                 size=12, color=NAVY)


def slide_team(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(prs, s, "Team & Hiring Plan",
                  "From 1 founder today → 16 FTE by Year 2.")
    rows = [
        ["Function", "Today", "Month 6", "Month 12", "Month 24"],
        ["Engineering", "1", "3", "5", "10"],
        ["Sales & demos", "0", "2", "5", "12"],
        ["Customer Success", "0", "1", "3", "6"],
        ["Marketing / Content", "0", "1", "2", "4"],
        ["Operations / Finance / HR", "0", "1", "1", "3"],
        ["TOTAL FTE", "1", "8", "16", "35"],
        ["Monthly burn (₹ L)", "—", "6.3", "14.5", "31.5"],
    ]
    _table(s, Inches(0.5), Inches(1.7), Inches(12.5), Inches(4.5), rows)
    _add_text(s, Inches(0.5), Inches(6.3), Inches(12.5), Inches(0.5),
              "Engineering hires the first quarter; sales scales fastest in months 6-18.",
              size=12, color=SLATE)


def slide_tech_scale(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(prs, s, "Technical Scalability",
                  "Cost-per-tenant asymptotes near ₹45/mo. Margin improves with scale.")
    rows = [
        ["Tenants", "Architecture", "Cost ₹/tenant/mo"],
        ["100", "Single droplet + Atlas M10", "₹130"],
        ["500", "+ Redis cache, CDN", "₹70"],
        ["2,000", "Atlas M20 + read replica", "₹47"],
        ["5,000", "Horizontal API, queue-based ingest, sharding", "₹45"],
        ["10,000", "Per-region clusters, multi-master", "₹45"],
    ]
    _table(s, Inches(0.5), Inches(1.7), Inches(12.5), Inches(3.5), rows)
    _add_rect(s, Inches(0.5), Inches(5.5), Inches(12.5), Inches(1.5), LIGHT, line=MUTED)
    _add_text(s, Inches(0.7), Inches(5.65), Inches(12), Inches(0.5),
              "Reliability targets",
              size=14, bold=True, color=NAVY)
    _add_text(s, Inches(0.7), Inches(6.05), Inches(12), Inches(0.4),
              "Year 1: 99.5% uptime · Year 2: 99.9% · RPO 1 hr · RTO 30 min",
              size=12, color=SLATE)


def slide_projections(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(prs, s, "24-Month Projections",
                  "1,000 tenants by M12 → 2,500 by M24 → ₹7.5 Cr ARR.")
    rows = [
        ["Month", "Cum. paying", "MRR (₹L)", "ARR (₹Cr)"],
        ["3", "95", "2.4", "0.3"],
        ["6", "280", "7.0", "0.8"],
        ["12", "1,000", "25.0", "3.0"],
        ["18", "1,750", "43.7", "5.2"],
        ["24", "2,500", "62.5", "7.5"],
    ]
    _table(s, Inches(0.5), Inches(1.7), Inches(7.5), Inches(3.6), rows)
    # Sensitivity sidebar
    _add_rect(s, Inches(8.5), Inches(1.7), Inches(4.5), Inches(3.6), LIGHT, line=MUTED)
    _add_text(s, Inches(8.7), Inches(1.85), Inches(4.2), Inches(0.4),
              "Sensitivity at M24", size=13, bold=True, color=NAVY)
    _add_text(s, Inches(8.7), Inches(2.4), Inches(4.2), Inches(0.4),
              "Bear · 0.6× exec", size=11, color=SLATE)
    _add_text(s, Inches(8.7), Inches(2.75), Inches(4.2), Inches(0.4),
              "₹4.5 Cr ARR · 1,500 tenants", size=14, bold=True, color=AMBER)
    _add_text(s, Inches(8.7), Inches(3.25), Inches(4.2), Inches(0.4),
              "Base · plan", size=11, color=SLATE)
    _add_text(s, Inches(8.7), Inches(3.6), Inches(4.2), Inches(0.4),
              "₹7.5 Cr ARR · 2,500 tenants", size=14, bold=True, color=BLUE)
    _add_text(s, Inches(8.7), Inches(4.1), Inches(4.2), Inches(0.4),
              "Bull · CA partner unlock", size=11, color=SLATE)
    _add_text(s, Inches(8.7), Inches(4.45), Inches(4.2), Inches(0.4),
              "₹10.5 Cr ARR · 3,500 tenants", size=14, bold=True, color=GREEN)
    # Path to default-alive
    _add_text(s, Inches(0.5), Inches(6.3), Inches(12.5), Inches(0.5),
              "EBITDA-positive: Month 30. Default-alive: Month 30 with current trajectory.",
              size=12, color=SLATE, bold=True)


def slide_2yr_plan(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(prs, s, "Two-Year Roadmap",
                  "Quarter-by-quarter execution.")
    rows = [
        ["Period", "Theme", "Key milestones", "Tenants"],
        ["Y1 Q1 (Mar–May)", "Foundation", "Atlas + DO + CI/CD + Sentry", "50"],
        ["Y1 Q2 (Jun–Aug)", "GTM ignition", "GST Portal, WhatsApp digest", "250"],
        ["Y1 Q3 (Sep–Nov)", "Salesman wedge", "Mobile-responsive, beat optimiser", "550"],
        ["Y1 Q4 (Dec–Feb)", "Year-1 close", "Mobile alpha, audit-log CSV", "1,000"],
        ["Y2 Q1 (Mar–May)", "Mobile + partner", "Mobile GA, CA programme open", "1,400"],
        ["Y2 Q2 (Jun–Aug)", "Geo expand", "Push to N. India + TN", "1,800"],
        ["Y2 Q3 (Sep–Nov)", "Enterprise tier", "SSO, API, white-label", "2,200"],
        ["Y2 Q4 (Dec–Feb)", "Series A", "₹15-20 Cr round closes", "2,500"],
    ]
    _table(s, Inches(0.5), Inches(1.7), Inches(12.5), Inches(5.0), rows, font_size=11)


def slide_risks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(prs, s, "Risks & Mitigations",
                  "What could go wrong, and what we're doing about it.")
    rows = [
        ["Risk", "Severity", "Likelihood", "Mitigation"],
        ["Tally Solutions launches competing dashboard", "High", "Medium",
         "We own dispatch + salesman they don't focus on"],
        ["Customer Tally version too old", "Med", "Med",
         "Already supports ERP 9 + Prime; agent degrades gracefully"],
        ["Meta WhatsApp policy tightens", "Med", "Med",
         "Diversify to SMS + email; user-initiated WA only"],
        ["Founder bus factor", "High", "Low",
         "Hire technical co-lead by M6; documented Production Playbook"],
        ["CAC inflation (Meta)", "Med", "High",
         "Diversify to YouTube + organic + CA partner"],
        ["GST policy change breaks sync", "Med", "Med",
         "Read-only consumer; we follow Tally's adaptation"],
        ["Customer data breach", "Catastrophic", "Low",
         "Atlas encryption, TLS, pen-test M4, cyber insurance"],
    ]
    _table(s, Inches(0.5), Inches(1.7), Inches(12.5), Inches(4.6), rows, font_size=10)


def slide_why_now(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(prs, s, "Why Now",
                  "Five tailwinds aligning in the same 24-month window.")
    items = [
        ("GST + e-invoicing mandate",
         "Forces every ₹5 Cr+ business to keep accurate Tally books — quality of source data is rising fast."),
        ("DPDP Act (Indian data privacy)",
         "Lift-and-shift cloud SaaS lose trust. We are data-residency-friendly: customer's books stay in India."),
        ("WhatsApp Business API maturing",
         "Zero-friction notifications now reliable. Was clunky 18 months ago."),
        ("Mobile-first owners in tier-2/3 cities",
         "~70% prefer phone over desktop for daily insights — perfect fit for our digest model."),
        ("Tally Solutions moving slowly",
         "Their 'Tally on AWS' has zero analytics. The wedge is wide open for 24-36 months."),
    ]
    y = Inches(1.8)
    for title, body in items:
        _add_rect(s, Inches(0.5), y, Inches(0.1), Inches(0.9), GREEN)
        _add_text(s, Inches(0.7), y, Inches(12), Inches(0.4),
                  title, size=14, bold=True, color=NAVY)
        _add_text(s, Inches(0.7), y + Inches(0.4), Inches(12.2), Inches(0.5),
                  body, size=11, color=SLATE)
        y += Inches(1.0)


def slide_ask(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(prs, s, "The Ask",
                  "₹4 Cr Seed at ₹16-20 Cr post-money.")
    # Big number on left
    _add_rect(s, Inches(0.5), Inches(1.9), Inches(5.5), Inches(4.8), NAVY)
    _add_text(s, Inches(0.5), Inches(2.2), Inches(5.5), Inches(0.6),
              "RAISING", size=14, color=RGBColor(0x93, 0xC5, 0xFD),
              align=PP_ALIGN.CENTER, bold=True)
    _add_text(s, Inches(0.5), Inches(2.9), Inches(5.5), Inches(1.5),
              "₹4 Cr", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.5), Inches(4.6), Inches(5.5), Inches(0.5),
              "at ₹16-20 Cr post-money",
              size=14, color=RGBColor(0xDB, 0xEA, 0xFE), align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.5), Inches(5.5), Inches(5.5), Inches(0.5),
              "24 months runway",
              size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.5), Inches(5.95), Inches(5.5), Inches(0.5),
              "Series A target: M24 at ₹7.5 Cr ARR",
              size=12, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
    # Use of funds on right
    rows = [
        ["Use of funds", "₹ Cr", "%"],
        ["Engineering hires (5 FTE × 24m)", "1.6", "40%"],
        ["Sales + CS hires (8 FTE × 24m)", "1.0", "25%"],
        ["Marketing / acquisition", "0.8", "20%"],
        ["Infra + tooling + compliance", "0.2", "5%"],
        ["Working capital buffer", "0.4", "10%"],
        ["TOTAL", "4.0", "100%"],
    ]
    _table(s, Inches(6.4), Inches(1.9), Inches(6.6), Inches(4.5), rows)
    _add_text(s, Inches(6.4), Inches(6.5), Inches(6.6), Inches(0.4),
              "Funded milestones: M12 → 1,000 tenants · M24 → 2,500 tenants · ₹7.5 Cr ARR",
              size=11, color=SLATE, bold=True)


def slide_close(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    _add_rect(s, 0, Inches(3.0), prs.slide_width, Inches(2.5), BLUE)
    _add_text(s, Inches(0.6), Inches(2.0), Inches(13), Inches(1.0),
              "We don't ask the customer to change anything.",
              size=32, bold=True, color=WHITE)
    _add_text(s, Inches(0.6), Inches(3.4), Inches(13), Inches(1.0),
              "We just turn on the lights.",
              size=32, bold=True, color=WHITE)
    _add_text(s, Inches(0.6), Inches(5.7), Inches(12), Inches(0.5),
              "Thank you.", size=20, color=RGBColor(0xDB, 0xEA, 0xFE),
              bold=True)
    _add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.4),
              "FLOWRA Insights · founders@flowra.in · Feb 2026",
              size=12, color=MUTED)


# ─── Visual chart slides (matplotlib-rendered PNGs) ────────────────────
def _chart_slide(prs, title, subtitle, png_path, caption=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(prs, s, title, subtitle)  # page_n auto-detected
    # Embed the PNG, centred with margins
    pic = s.shapes.add_picture(
        str(png_path),
        Inches(0.7), Inches(1.6),
        width=Inches(11.9),
    )
    # If picture exceeds vertical space, scale by height instead
    if pic.height > Inches(5.0):
        pic.width = int(pic.width * Inches(5.0) / pic.height)
        pic.height = Inches(5.0)
        pic.left = int((prs.slide_width - pic.width) / 2)
    if caption:
        _add_text(s, Inches(0.5), Inches(6.55), Inches(12.5), Inches(0.4),
                  caption, size=11, color=SLATE, bold=True,
                  align=PP_ALIGN.CENTER)


def slide_chart_funnel(prs, png):
    _chart_slide(prs,
                 "Acquisition Funnel",
                 "Per ₹50k Meta-ad spend cohort.",
                 png,
                 caption="3.5-month CAC payback at 78% gross margin.")


def slide_chart_growth(prs, png):
    _chart_slide(prs,
                 "Tenant Growth Curve",
                 "24-month plan with bear / base / bull sensitivity.",
                 png,
                 caption="Default-alive at M30. Series A target: M24 at ₹7.5 Cr ARR.")


def slide_chart_unit_econ(prs, png):
    _chart_slide(prs,
                 "Unit Economics",
                 "Why this scales profitably.",
                 png,
                 caption="LTV ₹85,000 · CAC ₹6,500 · Payback 3.5 mo · 78% gross margin.")


# ───────────────────────────────────────────────────────────────────────


def render_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Render chart PNGs first (deck embeds them)
    charts = render_all_charts()
    # Deck order — 19 slides total. Charts inserted where they tell the
    # story best:
    #   • Funnel chart sits right after the GTM table
    #   • Unit-economics chart sits after the pricing/unit-econ table
    #   • Growth chart sits before the projections table (visual lead-in)
    plan = [
        ("slide_cover",          lambda p: slide_cover(p)),
        ("slide_problem",        lambda p: slide_problem(p)),
        ("slide_solution",       lambda p: slide_solution(p)),
        ("slide_market",         lambda p: slide_market(p)),
        ("slide_competition",    lambda p: slide_competition(p)),
        ("slide_business_model", lambda p: slide_business_model(p)),
        # Visual: unit economics
        ("chart_unit_econ",      lambda p: slide_chart_unit_econ(p, charts["unit_econ"])),
        ("slide_gtm",            lambda p: slide_gtm(p)),
        # Visual: acquisition funnel
        ("chart_funnel",         lambda p: slide_chart_funnel(p, charts["funnel"])),
        ("slide_traction",       lambda p: slide_traction(p)),
        ("slide_team",           lambda p: slide_team(p)),
        ("slide_tech_scale",     lambda p: slide_tech_scale(p)),
        # Visual: tenant growth curve
        ("chart_growth",         lambda p: slide_chart_growth(p, charts["growth"])),
        ("slide_projections",    lambda p: slide_projections(p)),
        ("slide_2yr_plan",       lambda p: slide_2yr_plan(p)),
        ("slide_risks",          lambda p: slide_risks(p)),
        ("slide_why_now",        lambda p: slide_why_now(p)),
        ("slide_ask",            lambda p: slide_ask(p)),
        ("slide_close",          lambda p: slide_close(p)),
    ]
    for _name, fn in plan:
        fn(prs)

    prs.save(str(PPTX))
    print(f"Wrote {PPTX}  ({PPTX.stat().st_size // 1024} KB · {len(plan)} slides)")


def main():
    render_pdf()
    render_pptx()
    # Mirror to the frontend's public folder so the deployed app can serve
    # them at https://<host>/docs/<filename>
    public = Path("/app/frontend/public/docs")
    public.mkdir(parents=True, exist_ok=True)
    import shutil
    for f in (PDF, PPTX):
        shutil.copy2(f, public / f.name)
    # Production Playbook PDF (if it exists) — sibling artifact
    pp = DOCS / "PRODUCTION_PLAYBOOK.pdf"
    if pp.exists():
        shutil.copy2(pp, public / pp.name)
    print(f"Published copies → {public}")


if __name__ == "__main__":
    main()
